import io
import logging
from typing import Dict, List, Optional, Any, Union

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pymupdf
from langchain_core.tools import ToolException
from office365.runtime.auth.client_credential import ClientCredential
from office365.sharepoint.client_context import ClientContext
from pydantic import Field, PrivateAttr, create_model, model_validator
from transformers import BlipProcessor, BlipForConditionalGeneration

from .utils import read_docx_from_bytes
from ..elitea_base import BaseToolApiWrapper

# Constants
DEFAULT_LIST_LIMIT = 1000
DEFAULT_FILES_LIMIT = 100
SHARED_DOCUMENTS = "Shared Documents"
SUPPORTED_FILE_TYPES = ('.txt', '.docx', '.pdf', '.pptx')
BLIP_MODEL_NAME = "Salesforce/blip-image-captioning-base"

# Schema models for tool arguments
NoInput = create_model(
    "NoInput"
)

ReadList = create_model(
    "ReadList",
    list_title=(str, Field(description="Name of a Sharepoint list to be read.")),
    limit=(Optional[int], Field(description="Limit (maximum number) of list items to be returned", default=DEFAULT_LIST_LIMIT))
)

GetFiles = create_model(
    "GetFiles",
    folder_name=(Optional[str], Field(description="Folder name to get list of the files.", default=None)),
    limit_files=(Optional[int], Field(description="Limit (maximum number) of files to be returned. Can be called with synonyms, such as First, Top, etc., or can be reflected just by a number for example 'Top 10 files'. Use default value if not specified in a query WITH NO EXTRA CONFIRMATION FROM A USER", default=DEFAULT_FILES_LIMIT)),
)

ReadDocument = create_model(
    "ReadDocument",
    path=(str, Field(description="Contains the server-relative path of a document for reading.")),
    is_capture_image=(Optional[bool], Field(description="Determines if pictures in the document should be recognized.", default=False))
)


class SharepointApiWrapper(BaseToolApiWrapper):
    """
    A wrapper for SharePoint API operations.
    
    This class provides methods to interact with SharePoint lists, files, and documents.
    It supports authentication via client ID/secret or token.
    """
    site_url: str
    client_id: Optional[str] = None
    client_secret: Optional[str] = None
    token: Optional[str] = None
    _client: Optional[ClientContext] = PrivateAttr(default=None)  # Private attribute for the office365 client
    _image_processor: Optional[Any] = PrivateAttr(default=None)  # Lazy-loaded image processor
    _image_model: Optional[Any] = PrivateAttr(default=None)  # Lazy-loaded image model

    @model_validator(mode='before')
    @classmethod
    def validate_toolkit(cls, values):
        try:
            from office365.sharepoint.client_context import ClientContext
        except ImportError:
            raise ImportError(
                "`office365` package not found, please run "
               "`pip install office365-rest-python-client`"
            )

        site_url = values['site_url']
        client_id = values.get('client_id')
        client_secret = values.get('client_secret')
        token = values.get('token')
        _client = None

        if not ((client_id and client_secret) or token):
            raise ToolException("You have to define token or client id&secret.")

        try:
            if client_id and client_secret:
                credentials = ClientCredential(client_id, client_secret)
                _client = ClientContext(site_url).with_credentials(credentials)
                logging.info("SharePoint: Authenticated with client id/secret.")
            elif token:
                def _acquire_token():
                    return type('Token', (), {'tokenType': 'Bearer', 'accessToken': token})()
                _client = ClientContext(site_url).with_access_token(_acquire_token)
                logging.info("SharePoint: Authenticated with token.")

            values['_client'] = _client
            logging.info("SharePoint: Authentication successful and client assigned.")

        except Exception as e:
            logging.error(f"Failed to authenticate with SharePoint or create client: {str(e)}")
            values['_client'] = None

        return values


    def read_list(self, list_title: str, limit: int = DEFAULT_LIST_LIMIT) -> Union[List[Dict[str, Any]], ToolException]:
        """
        Reads a specified List in SharePoint site.
        
        Args:
            list_title: Name of the SharePoint list to read
            limit: Maximum number of list items to return (default is 1000)
            
        Returns:
            List of dictionaries containing list item properties or ToolException on error
        """
        if not self._client:
            logging.error("SharePoint client is not initialized")
            return ToolException("Cannot list items. SharePoint client is not initialized.")

        try:
            target_list = self._client.web.lists.get_by_title(list_title)
            self._client.load(target_list)
            self._client.execute_query()
            items = target_list.items.get().top(limit).execute_query()
            logging.info(f"{len(items)} items from SharePoint list '{list_title}' loaded successfully.")
            
            return [item.properties for item in items]
        except Exception as e:
            logging.error(f"Failed to load items from SharePoint list '{list_title}': {e}")
            return ToolException("Cannot list items. Please, double check List name and read permissions.")


    def get_files_list(self, folder_name: Optional[str] = None, limit_files: int = DEFAULT_FILES_LIMIT) -> Union[List[Dict[str, str]], ToolException]:
        """
        Lists files in a SharePoint folder.
        
        If folder name is specified, lists all files in this folder under Shared Documents path.
        If folder name is empty, lists all files under root catalog (Shared Documents).
        
        Args:
            folder_name: Name of the folder to list files from (default is None for root)
            limit_files: Maximum number of files to return (default is 100)
            
        Returns:
            List of dictionaries containing file properties or ToolException on error
        """
        if not self._client:
            logging.error("SharePoint client is not initialized")
            return ToolException("Cannot get files. SharePoint client is not initialized.")

        try:
            result = []
            target_folder_url = f"{SHARED_DOCUMENTS}/{folder_name}" if folder_name else SHARED_DOCUMENTS
            
            files = (self._client.web.get_folder_by_server_relative_path(target_folder_url)
                    .get_files(True)
                    .execute_query())

            for file in files:
                if len(result) >= limit_files:
                    break
                    
                result.append({
                    'Name': file.properties['Name'],
                    'Path': file.properties['ServerRelativeUrl'],
                    'Created': file.properties['TimeCreated'],
                    'Modified': file.properties['TimeLastModified'],
                    'Link': file.properties['LinkingUrl']
                })

            logging.info(f"Retrieved {len(result)} files from '{target_folder_url}'")
            return result
        except Exception as e:
            logging.error(f"Failed to load files from SharePoint folder '{folder_name}': {e}")
            return ToolException("Cannot get files. Please, double check folder name and read permissions.")

    def read_file(self, path: str, is_capture_image: bool = False) -> Union[str, ToolException]:
        """
        Reads file located at the specified server-relative path.
        
        Args:
            path: Server-relative path to the file
            is_capture_image: If True, attempts to describe images in PDF and PPTX files
            
        Returns:
            File content as string or ToolException on error
            
        Supports:
            - TXT: Plain text files
            - DOCX: Word documents
            - PDF: Adobe PDF files (with optional image description)
            - PPTX: PowerPoint presentations (with optional image description)
        """
        if not self._client:
            logging.error("SharePoint client is not initialized")
            return ToolException("File not found. SharePoint client is not initialized.")

        try:
            # Get file from SharePoint
            file = self._client.web.get_file_by_server_relative_path(path)
            self._client.load(file).execute_query()
            file_content = file.read()
            self._client.execute_query()
            
            # Check if file type is supported
            file_extension = self._get_file_extension(file.name)
            if file_extension not in SUPPORTED_FILE_TYPES:
                supported_types = ", ".join(ext[1:].upper() for ext in SUPPORTED_FILE_TYPES)
                return ToolException(f"Not supported type of file. Supported types are {supported_types} only.")
            
            # Process file based on its type
            if file_extension == '.txt':
                return self._process_txt_file(file_content)
            elif file_extension == '.docx':
                return self._process_docx_file(file_content)
            elif file_extension == '.pdf':
                return self._process_pdf_file(file_content, is_capture_image)
            elif file_extension == '.pptx':
                return self._process_pptx_file(file_content, is_capture_image)
            
        except Exception as e:
            logging.error(f"Failed to load file from SharePoint: {e}. Path: {path}")
            return ToolException("File not found. Please, check file name and path.")
    
    def _get_file_extension(self, filename: str) -> str:
        """Extract file extension from filename (lowercase)"""
        return filename[filename.rfind('.'):].lower() if '.' in filename else ''
    
    def _process_txt_file(self, file_content: bytes) -> Union[str, ToolException]:
        """Process TXT file content"""
        try:
            return file_content.decode('utf-8')
        except UnicodeDecodeError as e:
            logging.error(f"Error decoding file content: {e}")
            return ToolException("Error processing file content after download.")
    
    def _process_docx_file(self, file_content: bytes) -> str:
        """Process DOCX file content"""
        # read_docx_from_bytes already handles its errors and logs them
        return read_docx_from_bytes(file_content)
    
    def _process_pdf_file(self, file_content: bytes, is_capture_image: bool) -> str:
        """Process PDF file content with optional image description"""
        text_content = ''
        with pymupdf.open(stream=file_content, filetype="pdf") as pdf_doc:
            for page in pdf_doc:
                text_content += page.get_text()
                
                # Process images if requested
                if is_capture_image:
                    images = page.get_images(full=True)
                    for img in images:
                        xref = img[0]
                        base_image = pdf_doc.extract_image(xref)
                        img_bytes = base_image["image"]
                        
                        try:
                            pil_image = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                            caption = self.describe_image(pil_image)
                            text_content += caption
                        except Exception as img_err:
                            logging.warning(f"Could not process image in PDF: {img_err}")
                            text_content += "\n[Picture: processing error]\n"
        
        return text_content
    
    def _process_pptx_file(self, file_content: bytes, is_capture_image: bool) -> str:
        """Process PPTX file content with optional image description"""
        text_content = ''
        prs = Presentation(io.BytesIO(file_content))
        
        for slide in prs.slides:
            for shape in slide.shapes:
                # Extract text from text frames
                if shape.has_text_frame:
                    text_content += shape.text_frame.text + "\n"
                # Process images if requested
                elif is_capture_image and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        pil_image = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                        caption = self.describe_image(pil_image)
                        text_content += caption
                    except Exception as img_err:
                        logging.warning(f"Could not process image in PPTX: {img_err}")
                        text_content += "\n[Picture: processing error]\n"
        
        return text_content

    def describe_image(self, image: Image.Image) -> str:
        """
        Generate caption for an image using BLIP model.
        
        Args:
            image: PIL Image object to describe
            
        Returns:
            String containing the image description in format "[Picture: description]"
        """
        # Lazy-load the image processing models to save memory when not needed
        if self._image_processor is None or self._image_model is None:
            try:
                self._image_processor = BlipProcessor.from_pretrained(BLIP_MODEL_NAME)
                self._image_model = BlipForConditionalGeneration.from_pretrained(BLIP_MODEL_NAME)
                logging.info("Image description models loaded successfully")
            except Exception as e:
                logging.error(f"Failed to load image description models: {e}")
                return "\n[Picture: description unavailable - model loading failed]\n"
        
        try:
            inputs = self._image_processor(image, return_tensors="pt")
            out = self._image_model.generate(**inputs)
            caption = self._image_processor.decode(out[0], skip_special_tokens=True)
            return f"\n[Picture: {caption}]\n"
        except Exception as e:
            logging.error(f"Failed to describe image: {e}")
            return "\n[Picture: description unavailable]\n"


    def get_available_tools(self) -> List[Dict[str, Any]]:
        """
        Returns the list of available tools provided by this wrapper.
        
        Returns:
            List of tool definitions with name, description, schema and reference
        """
        return [
            {
                "name": "read_list",
                "description": self.read_list.__doc__,
                "args_schema": ReadList,
                "ref": self.read_list
            },
            {
                "name": "get_files_list",
                "description": self.get_files_list.__doc__,
                "args_schema": GetFiles,
                "ref": self.get_files_list
            },
            {
                "name": "read_document",
                "description": self.read_file.__doc__,
                "args_schema": ReadDocument,
                "ref": self.read_file
            }
        ]
